import datetime
import tempfile
import os
import shutil
from pathlib import Path
import pymupdf as fitz

from PIL import Image
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF
from odf.opendocument import OpenDocumentText
from odf.text import P

from rest_framework.views import APIView
from rest_framework.permissions import AllowAny
from rest_framework.response import Response
from rest_framework import status
from django.http import HttpResponse
import django.conf

from dropfiles.models import FileUpload


class ConvertView(APIView):
    permission_classes = [AllowAny]

    def post(self, request):
        try:
            data = request.data
            uploaded_file = data.get('file')
            from_format = data.get('from_format', '').lower()
            to_format = data.get('to_format', '').lower()
            
            if not uploaded_file:
                return Response({'error': 'Файл не загружен'}, status=status.HTTP_400_BAD_REQUEST)
            
            with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{from_format}') as tmp_file:
                tmp_file.write(uploaded_file.read())
                input_path = tmp_file.name
            
            output_path = str(Path(input_path).with_suffix(f'.{to_format}'))

            if not self._is_conversion_supported(from_format, to_format):
                os.unlink(input_path)
                return Response({'error': f'Конвертация из {from_format.upper()} в {to_format.upper()} недоступна'}, status=status.HTTP_400_BAD_REQUEST)

            success, message = self._convert_file(input_path, output_path, from_format, to_format)
            if not success:
                if os.path.exists(input_path): os.unlink(input_path)
                if os.path.exists(output_path): os.unlink(output_path)
                return Response({'error': f'Ошибка конвертации: {message}'}, status=status.HTTP_400_BAD_REQUEST)
            
            with open(output_path, 'rb') as file:
                converted_data = file.read()
            
            original_size = uploaded_file.size
            converted_size = len(converted_data)
            file_name = f'{Path(uploaded_file.name).stem}.{to_format}'

            os.unlink(input_path)
            os.unlink(output_path)

            response = HttpResponse(converted_data, content_type='application/octet-stream')
            response['Content-Disposition'] = f"attachment; filename='{file_name}'"
            response['X-Original-Size'] = str(original_size)
            response['X-Converted-Size'] = str(converted_size)
            response['X-From-Format'] = from_format.upper()
            response['X-To-Format'] = to_format.upper()
            response['Access-Control-Expose-Headers'] = 'X-Original-Size, X-Converted-Size, X-From-Format, X-To-Format, Content-Disposition'
            
            user = request.user if request.user.is_authenticated else None

            FileUpload.objects.create(
                files=[file_name],
                size=converted_size,
                created_at=datetime.datetime.now(),
                user=user
            )
            return response
            
        except Exception as e:
            try:
                if 'input_path' in locals() and os.path.exists(input_path): os.unlink(input_path)
                if 'output_path' in locals() and os.path.exists(output_path): os.unlink(output_path)
            except:
                pass
                
            return Response(
                {'error': f'Ошибка при обработке: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

    def _is_conversion_supported(self, from_format, to_format):
        supported = {
            'PDF': ['DOCX', 'TXT', 'JPG', 'PNG', 'PPTX', 'JPEG'],
            'DOCX': ['PDF', 'TXT', 'ODT'],
            'JPG': ['PNG', 'PDF', 'WEBP', 'JPEG'],
            'JPEG': ['PNG', 'PDF', 'WEBP', 'JPG'],
            'PNG': ['JPG', 'PDF', 'WEBP', 'JPEG'],
            'TXT': ['PDF', 'DOCX', 'CSV'],
            'PPTX': ['PDF', 'JPG', 'PNG'],
            'CSV': ['TXT'],
            'WEBP': ['JPG', 'PNG', 'PDF', 'JPEG'],
            'ODT': ['DOCX', 'PDF', 'TXT'],
        }
        return from_format.upper() in supported and to_format.upper() in supported.get(from_format.upper(), [])
    
    def _convert_file(self, input_path, output_path, from_format, to_format):
        try:
            if from_format in ['jpg', 'jpeg', 'png', 'webp'] and to_format in ['jpg', 'jpeg', 'png', 'webp']:
                return self._convert_image(input_path, output_path, to_format)
            
            elif from_format == 'pdf':
                if to_format == 'txt':
                    return self._pdf_to_text(input_path, output_path)
                elif to_format in ['jpg', 'jpeg', 'png']:
                    return self._pdf_to_image(input_path, output_path, to_format)
                elif to_format == 'docx':
                    return self._pdf_to_docx(input_path, output_path)
                elif to_format == 'pptx':
                    return self._pdf_to_pptx(input_path, output_path)

            elif from_format == 'docx':
                if to_format == 'pdf':
                    return self._docx_to_pdf(input_path, output_path)
                elif to_format == 'txt':
                    return self._docx_to_txt(input_path, output_path)
                elif to_format == 'odt':
                    return self._docx_to_odt(input_path, output_path)

            elif from_format == 'txt':
                if to_format == 'docx':
                    return self._txt_to_docx(input_path, output_path)
                elif to_format == 'pdf':
                    return self._txt_to_pdf(input_path, output_path)
                elif to_format == 'csv':
                    shutil.copy(input_path, output_path)
                    return True, 'Успешно'
            elif from_format == 'csv' and to_format == 'txt':
                shutil.copy(input_path, output_path)
                return True, 'Успешно'

            elif from_format == 'pptx':
                return self._pptx_to_layouts(input_path, output_path, to_format)

            elif from_format == 'odt':
                return self._odt_to_layouts(input_path, output_path, to_format)
                
            return False, f'Конвертация {from_format.upper()} -> {to_format.upper()} не поддерживается'
                
        except Exception as e:
            return False, str(e)
    
    def _convert_image(self, input_path, output_path, to_format):
        try:
            img = Image.open(input_path)
            pil_format = to_format.upper()
            
            if pil_format == 'JPG':
                pil_format = 'JPEG'
                
            if pil_format == 'JPEG':
                if img.mode in ('RGBA', 'LA', 'P'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    background.paste(img, mask=img.split()[3] if len(img.split()) == 4 else None)
                    img = background
            
            img.save(output_path, format=pil_format)
            return True, 'Изображение сконвертировано'
        except Exception as error:
            return False, str(error)
        
    def _pdf_to_text(self, input_path, output_path):
        try:
            doc = fitz.open(input_path)
            text = ''.join([page.get_text() for page in doc])
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            return True, 'Текст из PDF извлечен'
        except Exception as error:
            return False, str(error)

    def _pdf_to_image(self, input_path, output_path, to_format):
        try:
            doc = fitz.open(input_path)
            page = doc[0]
            pix = page.get_pixmap()
            pil_format = 'JPEG' if to_format.lower() in ['jpg', 'jpeg'] else 'PNG'
            
            img = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
            img.save(output_path, format=pil_format)
            return True, 'Страница PDF сохранена как изображение'
        except Exception as e:
            return False, str(e)

    def _pdf_to_docx(self, input_path, output_path):
        try:
            from pdf2docx import Converter
            cv = Converter(input_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
            return True, 'PDF успешно конвертирован в DOCX'
        except Exception as e:
            return False, str(e)

    def _pdf_to_pptx(self, input_path, output_path):
        try:
            doc = fitz.open(input_path)
            prs = Presentation()
            for page in doc:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                text = page.get_text("text")
                
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                tf = txBox.text_frame
                tf.text = text
                
            prs.save(output_path)
            return True, 'PDF конвертирован в PPTX (текст)'
        except Exception as e:
            return False, str(e)

    def _docx_to_pdf(self, input_path, output_path):
        try:
            doc = Document(input_path)
            pdf = FPDF()
            pdf.add_page()
            
            font_path = os.path.join(django.conf.settings.MEDIA_ROOT, 'Arial.ttf')
            
            pdf.add_font('Arial', '', font_path, uni=True)
            pdf.set_font('Arial', size=11) 
                
            for p in doc.paragraphs:
                if p.text.strip():
                    pdf.multi_cell(0, 7, txt=p.text)
                    pdf.ln(2)
            pdf.output(output_path)
            return True, "Успешно"
        except Exception as e:
            return False, str(e)

    def _docx_to_txt(self, input_path, output_path):
        try:
            doc = Document(input_path)
            text = '\n'.join([p.text for p in doc.paragraphs])
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            return True, 'Успешно'
        except Exception as e:
            return False, str(e)

    def _docx_to_odt(self, input_path, output_path):
        try:
            doc = Document(input_path)
            odt_doc = OpenDocumentText()
            for p in doc.paragraphs:
                odt_doc.text.addElement(P(text=p.text))
            odt_doc.save(output_path)
            return True, 'Успешно'
        except Exception as e:
            return False, str(e)

    def _txt_to_docx(self, input_path, output_path):
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                lines = f.readlines()
            doc = Document()
            for line in lines:
                doc.add_paragraph(line.strip('\n'))
            doc.save(output_path)
            return True, 'Успешно'
        except Exception as e:
            return False, str(e)

    def _txt_to_pdf(self, input_path, output_path):
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                text = f.read()
            pdf = FPDF()
            pdf.add_page()
            font_path = os.path.join(django.conf.settings.MEDIA_ROOT, 'Arial.ttf')
            try:
                pdf.add_font('Arial', '', font_path, uni=True)
                pdf.set_font('Arial', size=11)
            except:
                pass
                
            pdf.multi_cell(0, 6, txt=text)
            pdf.output(output_path)
            return True, 'Успешно'
        except Exception as e:
            return False, str(e)

    def _pptx_to_layouts(self, input_path, output_path, to_format):
        try:
            prs = Presentation(input_path)
            
            if to_format in ['jpg', 'jpeg', 'png']:
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.shape_type == 13:
                            with open(output_path, 'wb') as f:
                                f.write(shape.image.blob)
                            return True, 'Первое найденное изображение успешно извлечено'
                return False, 'В презентации не найдено изображений для извлечения'

            text_content = []
            for i, slide in enumerate(prs.slides):
                text_content.append(f'--- Слайд {i+1} ---')
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text_content.append(shape.text)
            
            full_text = '\n'.join(text_content)

            if to_format == 'txt':
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(full_text)
                return True, 'Текст презентации сохранен в TXT'

            if to_format == 'pdf':
                pdf = FPDF()
                pdf.add_page()
                font_path = os.path.join(django.conf.settings.MEDIA_ROOT, 'Arial.ttf')
                try:
                    pdf.add_font('Arial', '', font_path, uni=True)
                    pdf.set_font('Arial', size=11)
                except:
                    pass
                
                pdf.multi_cell(0, 6, txt=full_text)
                pdf.output(output_path)
                return True, 'Текст презентации сохранен в PDF'

            return False, f'Формат {to_format.upper()} не поддерживается для PPTX'

        except Exception as e:
            return False, str(e)

    def _odt_to_layouts(self, input_path, output_path, to_format):
        try:
            from odf import text, teletype
            from odf.opendocument import load
            odt_doc = load(input_path)
            paragraphs = odt_doc.getElementsByType(text.P)
            extracted_text = '\n'.join([teletype.extractText(p) for p in paragraphs])
            
            if to_format == 'txt':
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(extracted_text)
                return True, 'Успешно'
            elif to_format == 'docx':
                doc = Document()
                for p in extracted_text.split('\n'):
                    doc.add_paragraph(p)
                doc.save(output_path)
                return True, 'Успешно'
            elif to_format == 'pdf':
                pdf = FPDF()
                pdf.add_page()
                font_path = os.path.join(django.conf.settings.MEDIA_ROOT, 'Arial.ttf')
                try:
                    pdf.add_font('Arial', '', font_path, uni=True)
                    pdf.set_font('Arial', size=11)
                except:
                    pass
                pdf.multi_cell(0, 6, txt=extracted_text)
                pdf.output(output_path)
                return True, 'Успешно'
            return False, 'Неподдерживаемый целевой формат для ODT'
        except Exception as e:
            return False, str(e)